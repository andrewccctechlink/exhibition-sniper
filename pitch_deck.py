"""Pitch deck generator – creates a Morandi-themed PowerPoint presentation."""

import json
import logging
import os
from datetime import datetime

logger = logging.getLogger(__name__)


class PitchDeckGenerator:
    """Generate pitch deck PowerPoint files."""

    def __init__(self, storage_path):
        self.storage_path = storage_path
        os.makedirs(storage_path, exist_ok=True)

    def generate(self, company_name, products, contact_info=None):
        """Generate a pitch deck. Returns dict with filepath and download_url."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            # Fallback: generate a JSON summary instead
            return self._generate_json_fallback(company_name, products, contact_info)

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Morandi color palette
        bg_color = RGBColor(0xF5, 0xF0, 0xEB)  # warm beige
        title_color = RGBColor(0x4A, 0x4A, 0x4A)  # dark grey
        accent_color = RGBColor(0x8B, 0x7D, 0x6B)  # warm brown

        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_color

        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Product Solutions for {company_name}"
        p.font.size = Pt(36)
        p.font.color.rgb = title_color
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Product slides
        for product in products[:10]:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = bg_color

            name = product.get("name", "Product")
            # Title
            tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
            p = tb.text_frame.paragraphs[0]
            p.text = name
            p.font.size = Pt(28)
            p.font.color.rgb = title_color
            p.font.bold = True

            # Details
            details_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11), Inches(5))
            tf = details_box.text_frame
            tf.word_wrap = True

            if product.get("headline"):
                p = tf.add_paragraph()
                p.text = product["headline"]
                p.font.size = Pt(18)
                p.font.color.rgb = accent_color
                p.font.italic = True

            for point in product.get("key_selling_points", []):
                p = tf.add_paragraph()
                p.text = f"• {point}"
                p.font.size = Pt(14)
                p.font.color.rgb = title_color

            if product.get("price"):
                p = tf.add_paragraph()
                p.text = f"\nPrice: {product['price']}"
                p.font.size = Pt(14)
                p.font.color.rgb = accent_color

        # Save
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_name = "".join(c if c.isalnum() or c in " _-" else "" for c in company_name)[:30]
        filename = f"PitchDeck_{safe_name}_{timestamp}.pptx"
        filepath = os.path.join(self.storage_path, filename)
        prs.save(filepath)

        return {
            "filepath": filepath,
            "filename": filename,
            "download_url": f"/api/download/{filename}",
            "slides": 1 + len(products[:10]),
        }

    def _generate_json_fallback(self, company_name, products, contact_info):
        """Fallback when python-pptx is not available."""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_name = "".join(c if c.isalnum() or c in " _-" else "" for c in company_name)[:30]
        filename = f"PitchDeck_{safe_name}_{timestamp}.json"
        filepath = os.path.join(self.storage_path, filename)

        deck = {
            "title": f"Product Solutions for {company_name}",
            "generated_at": datetime.now().isoformat(),
            "products": products,
            "contact_info": contact_info,
        }
        with open(filepath, "w", encoding="utf-8") as f:
            json.dump(deck, f, ensure_ascii=False, indent=2)

        return {
            "filepath": filepath,
            "filename": filename,
            "download_url": f"/api/download/{filename}",
            "slides": len(products),
        }
